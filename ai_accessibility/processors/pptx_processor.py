"""PowerPoint accessibility processor."""

import io
from typing import Optional
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from .base import BaseProcessor
from utils.accessibility import AccessibilityIssue, Severity
from utils.claude_client import ClaudeClient


class PowerPointProcessor(BaseProcessor):
    """Processor for PowerPoint documents."""

    # Generic/non-descriptive link texts to fix
    GENERIC_LINK_TEXTS = {
        'here', 'click here', 'click', 'link', 'this link',
        'read more', 'learn more', 'more', 'info', 'details'
    }

    def __init__(self, claude_client: Optional[ClaudeClient] = None):
        super().__init__(claude_client)

    def get_file_extension(self) -> str:
        return ".pptx"

    def process(self, content: bytes, filename: str = "") -> bytes:
        """
        Process PowerPoint for WCAG 2.1 AA accessibility.

        Args:
            content: PPTX content as bytes
            filename: Original filename

        Returns:
            Accessible PPTX as bytes
        """
        self.reset_report()

        # Open presentation
        prs = Presentation(io.BytesIO(content))

        # Apply accessibility fixes
        self._check_slide_titles(prs)
        self._add_alt_text_to_images(prs)
        self._check_reading_order(prs)
        self._check_table_headers(prs)
        self._add_slide_notes_for_complex_content(prs)
        self._check_hyperlinks(prs)
        self._check_color_only_information(prs)
        self._check_ambiguous_references(prs)

        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    def _check_slide_titles(self, prs: Presentation):
        """Ensure all slides have titles."""
        for slide_num, slide in enumerate(prs.slides, 1):
            has_title = False
            title_text = ""

            # Check for title placeholder
            if slide.shapes.title:
                title_shape = slide.shapes.title
                if title_shape.has_text_frame:
                    title_text = title_shape.text_frame.text.strip()
                    has_title = bool(title_text)

            # If no title placeholder, check for any shape that might be a title
            if not has_title:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        # Check if it looks like a title (short, at top of slide)
                        if text and len(text) < 100:
                            if hasattr(shape, 'top') and shape.top < Inches(1.5):
                                title_text = text
                                has_title = True
                                break

            if not has_title:
                self.report.add_issue(AccessibilityIssue(
                    wcag_criterion="2.4.6",
                    severity=Severity.ERROR,
                    description=f"Slide {slide_num} has no title",
                    suggestion="Add a descriptive title to help navigation"
                ))

                # Try to generate a title based on content
                slide_text = self._get_slide_text(slide)[:500]
                if slide_text:
                    try:
                        # We don't modify the slide structure, just report
                        self.report.add_warning(
                            f"Slide {slide_num} needs a title. Content preview: '{slide_text[:100]}...'"
                        )
                    except Exception:
                        pass

    def _add_alt_text_to_images(self, prs: Presentation):
        """Add alt text to images and shapes."""
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = self._get_slide_text(slide)[:300]

            for shape in slide.shapes:
                # Check if shape is an image or picture
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    alt_text = self._get_alt_text(shape)

                    if not alt_text:
                        self.report.add_issue(AccessibilityIssue(
                            wcag_criterion="1.1.1",
                            severity=Severity.ERROR,
                            description=f"Image on slide {slide_num} missing alt text",
                            location=f"Shape: {shape.name}"
                        ))

                        # Generate alt text
                        try:
                            # Extract image data
                            image = shape.image
                            image_bytes = image.blob
                            content_type = image.content_type

                            # Generate description
                            new_alt = self.claude_client.generate_alt_text(
                                image_data=image_bytes,
                                media_type=content_type,
                                surrounding_text=slide_text
                            )

                            if new_alt == "DECORATIVE":
                                self._set_alt_text(shape, "", decorative=True)
                                self.report.add_fix(f"Marked image as decorative on slide {slide_num}")
                            else:
                                self._set_alt_text(shape, new_alt)
                                self.report.add_fix(f"Added alt text to image on slide {slide_num}")

                        except Exception as e:
                            self.report.add_warning(f"Could not generate alt text for slide {slide_num}: {str(e)}")

                # Check other shape types that might need descriptions
                elif shape.shape_type in [MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.DIAGRAM,
                                          MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT]:
                    alt_text = self._get_alt_text(shape)

                    if not alt_text:
                        self.report.add_issue(AccessibilityIssue(
                            wcag_criterion="1.1.1",
                            severity=Severity.WARNING,
                            description=f"Complex shape on slide {slide_num} may need description",
                            location=f"Shape type: {shape.shape_type}, Name: {shape.name}"
                        ))

    def _get_alt_text(self, shape) -> str:
        """Get alt text from a shape."""
        try:
            # Access the descr attribute from the shape's XML
            nvSpPr = shape._element.find(qn('p:nvSpPr')) or shape._element.find(qn('p:nvPicPr'))
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(qn('p:cNvPr'))
                if cNvPr is not None:
                    return cNvPr.get('descr', '')

            # Alternative: check nvPicPr for pictures
            nvPicPr = shape._element.find('.//' + qn('p:nvPicPr'))
            if nvPicPr is not None:
                cNvPr = nvPicPr.find(qn('p:cNvPr'))
                if cNvPr is not None:
                    return cNvPr.get('descr', '')

        except Exception:
            pass

        return ''

    def _set_alt_text(self, shape, alt_text: str, decorative: bool = False):
        """Set alt text for a shape."""
        try:
            # Find the cNvPr element
            nvPicPr = shape._element.find('.//' + qn('p:nvPicPr'))
            if nvPicPr is not None:
                cNvPr = nvPicPr.find(qn('p:cNvPr'))
                if cNvPr is not None:
                    cNvPr.set('descr', alt_text)
                    if decorative:
                        # For decorative images, also set title to empty
                        cNvPr.set('title', '')

            nvSpPr = shape._element.find('.//' + qn('p:nvSpPr'))
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(qn('p:cNvPr'))
                if cNvPr is not None:
                    cNvPr.set('descr', alt_text)

        except Exception as e:
            self.report.add_warning(f"Could not set alt text: {str(e)}")

    def _check_reading_order(self, prs: Presentation):
        """Check reading order of slide elements."""
        for slide_num, slide in enumerate(prs.slides, 1):
            shapes = list(slide.shapes)

            if len(shapes) > 3:
                # Check if shapes are ordered logically (top to bottom, left to right)
                positioned_shapes = []
                for shape in shapes:
                    if hasattr(shape, 'top') and hasattr(shape, 'left'):
                        positioned_shapes.append((shape.top, shape.left, shape.name))

                if positioned_shapes:
                    # Sort by expected reading order
                    expected_order = sorted(positioned_shapes, key=lambda x: (x[0], x[1]))
                    actual_order = positioned_shapes

                    # Check how different they are
                    if actual_order != expected_order:
                        self.report.add_issue(AccessibilityIssue(
                            wcag_criterion="1.3.2",
                            severity=Severity.WARNING,
                            description=f"Slide {slide_num} reading order may not match visual order",
                            suggestion="Review reading order in Selection Pane"
                        ))

    def _check_table_headers(self, prs: Presentation):
        """Check tables for header rows."""
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table

                    # Check if first row is marked as header
                    # In python-pptx, we check the firstRow property
                    has_header = False
                    try:
                        # Check table properties
                        tbl = shape._element.find('.//' + qn('a:tbl'))
                        if tbl is not None:
                            tblPr = tbl.find(qn('a:tblPr'))
                            if tblPr is not None:
                                has_header = tblPr.get('firstRow') == '1'
                    except Exception:
                        pass

                    if not has_header:
                        self.report.add_issue(AccessibilityIssue(
                            wcag_criterion="1.3.1",
                            severity=Severity.WARNING,
                            description=f"Table on slide {slide_num} may not have header row marked",
                            suggestion="Enable 'Header Row' in Table Design options"
                        ))

                    # Generate caption if table is complex
                    if table.rows and len(table.rows) > 3:
                        self.report.add_issue(AccessibilityIssue(
                            wcag_criterion="1.3.1",
                            severity=Severity.INFO,
                            description=f"Complex table on slide {slide_num}",
                            suggestion="Consider adding table description in slide notes"
                        ))

    def _add_slide_notes_for_complex_content(self, prs: Presentation):
        """Add notes for slides with complex visual content."""
        for slide_num, slide in enumerate(prs.slides, 1):
            has_complex_content = False
            complex_shapes = []

            for shape in slide.shapes:
                if shape.shape_type in [MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.DIAGRAM,
                                        MSO_SHAPE_TYPE.GROUP]:
                    has_complex_content = True
                    complex_shapes.append(shape.shape_type)

            if has_complex_content:
                # Check if notes already exist
                notes_slide = slide.notes_slide if slide.has_notes_slide else None
                existing_notes = ""

                if notes_slide:
                    notes_frame = notes_slide.notes_text_frame
                    existing_notes = notes_frame.text if notes_frame else ""

                if not existing_notes or len(existing_notes) < 50:
                    self.report.add_issue(AccessibilityIssue(
                        wcag_criterion="1.1.1",
                        severity=Severity.WARNING,
                        description=f"Slide {slide_num} has complex visuals but minimal speaker notes",
                        suggestion="Add detailed descriptions in speaker notes for accessibility"
                    ))

    def _check_hyperlinks(self, prs: Presentation):
        """Check hyperlinks for accessibility."""
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.hyperlink and run.hyperlink.address:
                                link_text = run.text
                                href = run.hyperlink.address

                                # Check for non-descriptive link text
                                from utils.accessibility import AccessibilityChecker
                                issue = AccessibilityChecker.check_link_text(link_text, href)

                                if issue:
                                    issue.location = f"Slide {slide_num}"
                                    self.report.add_issue(issue)

    def _get_slide_text(self, slide) -> str:
        """Get all text from a slide."""
        text_parts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text_parts.append(paragraph.text)

        return ' '.join(text_parts)

    def _check_color_only_information(self, prs: Presentation):
        """Check for color-only information in slides."""
        import re

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = self._get_slide_text(slide)

            # Check for color-only references in text
            color_ref_pattern = r'(the|see|marked in|shown in|highlighted in|in)\s+(red|green|blue|yellow|orange|purple|pink)\s*(text|items?|sections?|areas?|cells?)?'

            for match in re.finditer(color_ref_pattern, slide_text, re.IGNORECASE):
                self.report.add_issue(AccessibilityIssue(
                    wcag_criterion="1.4.1",
                    severity=Severity.ERROR,
                    description=f"Color-only reference on slide {slide_num}: '{match.group(0)}'",
                    suggestion="Add non-color indicator (bold, underline, symbols) in addition to color"
                ))

            # Check shapes for color-only formatting
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Check if run has color but no other distinguishing format
                            try:
                                # Check if color exists and has rgb property
                                if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                    # Get RGB as hex
                                    color = run.font.color.rgb
                                    # Check if it's a non-black/white color (indicating emphasis)
                                    if color and str(color) not in ['000000', 'FFFFFF', '000', 'FFF']:
                                        # Check if text has other formatting
                                        has_bold = run.font.bold
                                        has_italic = run.font.italic
                                        has_underline = run.font.underline

                                        if not has_bold and not has_italic and not has_underline:
                                            text = run.text[:30] if len(run.text) > 30 else run.text
                                            if text.strip():
                                                self.report.add_issue(AccessibilityIssue(
                                                    wcag_criterion="1.4.1",
                                                    severity=Severity.WARNING,
                                                    description=f"Color-only emphasis on slide {slide_num}: '{text}'",
                                                    suggestion="Add bold, italic, or underline in addition to color"
                                                ))
                            except Exception:
                                pass

    def _check_ambiguous_references(self, prs: Presentation):
        """Check for ambiguous page/visual-only references in slides."""
        import re

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = self._get_slide_text(slide)

            # Pattern for "see slide X" or "on slide X" (slide-specific)
            slide_ref_pattern = r'([Ss]ee|[Oo]n|[Rr]efer to)\s+slide\s+(\d+)'
            for match in re.finditer(slide_ref_pattern, slide_text):
                self.report.add_issue(AccessibilityIssue(
                    wcag_criterion="1.3.1",
                    severity=Severity.INFO,
                    description=f"Slide reference on slide {slide_num}: '{match.group(0)}'",
                    suggestion="Consider using slide titles for clearer navigation"
                ))

            # Pattern for "see page X" or "on page X"
            page_ref_pattern = r'([Ss]ee|[Oo]n|[Rr]efer to)\s+page\s+(\d+)'
            for match in re.finditer(page_ref_pattern, slide_text):
                self.report.add_issue(AccessibilityIssue(
                    wcag_criterion="1.3.1",
                    severity=Severity.WARNING,
                    description=f"Ambiguous page reference on slide {slide_num}: '{match.group(0)}'",
                    suggestion="Use descriptive references instead of page numbers"
                ))

            # Pattern for "above" or "below" references
            position_ref_pattern = r'(the|see|as shown)\s+(above|below|following|previous)\s+(figure|table|section|image|diagram|slide|chart)'
            for match in re.finditer(position_ref_pattern, slide_text, re.IGNORECASE):
                self.report.add_issue(AccessibilityIssue(
                    wcag_criterion="1.3.1",
                    severity=Severity.INFO,
                    description=f"Position-based reference on slide {slide_num}: '{match.group(0)}'",
                    suggestion="Use explicit labels or names for referenced content"
                ))

            # Pattern for "on the left/right" references
            spatial_ref_pattern = r'(on the|to the|at the)\s+(left|right|top|bottom|center|corner)'
            for match in re.finditer(spatial_ref_pattern, slide_text, re.IGNORECASE):
                self.report.add_issue(AccessibilityIssue(
                    wcag_criterion="1.3.1",
                    severity=Severity.INFO,
                    description=f"Spatial reference on slide {slide_num}: '{match.group(0)}'",
                    suggestion="Provide context that doesn't rely on visual position"
                ))
